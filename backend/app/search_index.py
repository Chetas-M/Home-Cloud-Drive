"""
Home Cloud Drive - Rich Search Indexing
========================================
Extracts searchable text from a wide variety of file formats so that
PDFs, Office documents, and (optionally) images become discoverable by
content.

Extractor matrix
----------------
Format              Library             Always available?
------------------  ------------------  -----------------
Plain text / code   built-in            ✓
PDF                 pypdf               optional (graceful fallback)
                    pdfminer.six        optional (higher fidelity)
Word (.docx)        python-docx         optional
Excel (.xlsx/.xls)  openpyxl / xlrd     optional
PowerPoint (.pptx)  python-pptx         optional
Images (OCR)        pytesseract+Pillow  optional (requires Tesseract binary)

When an optional library is missing the extractor silently falls back
to returning ``None``; the file simply won't have content-indexed text.
Install optional extras with:

    pip install pypdf pdfminer.six python-docx openpyxl python-pptx
    # For OCR (requires Tesseract system package as well):
    pip install pytesseract

Snippet helpers
---------------
``build_match_context`` builds a short highlighted snippet for search
result cards.  ``build_search_document`` is the top-level entry point
called from both synchronous indexing paths and the background job
queue.
"""
from __future__ import annotations

import io
import logging
import os
import re
from typing import Iterable, List, Optional

from app.models import File as FileModel

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

MAX_INDEX_BYTES = 512 * 1024          # 512 KB of raw bytes to read from file
MAX_INDEX_CHARS = 200_000             # Upper bound on indexed text chars
MAX_MATCH_CONTEXT_CHARS = 200

TEXT_FILE_TYPES = {"text"}
TEXT_EXTENSIONS = {
    "txt", "md", "markdown", "json", "xml", "html", "htm", "css", "js", "jsx",
    "ts", "tsx", "py", "java", "cpp", "c", "h", "hpp", "csv", "log", "yaml",
    "yml", "ini", "toml", "env", "sql", "sh", "bash", "zsh", "rb", "go",
    "rs", "php", "swift", "kt", "scala", "r", "m", "vb", "cs", "dart",
}
PDF_EXTENSIONS = {"pdf"}
WORD_EXTENSIONS = {"docx"}
EXCEL_EXTENSIONS = {"xlsx", "xlsm", "xls"}
PPTX_EXTENSIONS = {"pptx"}
IMAGE_EXTENSIONS = {"jpg", "jpeg", "png", "bmp", "tiff", "tif", "gif", "webp"}

# ---------------------------------------------------------------------------
# Utility helpers
# ---------------------------------------------------------------------------


def normalize_whitespace(value: str) -> str:
    return re.sub(r"\s+", " ", value or "").strip()


def _extension(filename: str) -> str:
    return os.path.splitext(filename or "")[1].lower().lstrip(".")


def _truncate(text: str) -> str:
    return text[:MAX_INDEX_CHARS]


# ---------------------------------------------------------------------------
# Individual extractors
# ---------------------------------------------------------------------------


def _extract_plain_text(path: str) -> Optional[str]:
    """Read up to MAX_INDEX_BYTES of a text / source-code file."""
    try:
        with open(path, "rb") as fh:
            raw = fh.read(MAX_INDEX_BYTES)
        text = raw.decode("utf-8", errors="ignore")
        return _truncate(normalize_whitespace(text)) or None
    except OSError as exc:
        logger.debug("search_index: plain-text read failed for %s: %s", path, exc)
        return None


def _extract_pdf(path: str) -> Optional[str]:
    """
    Extract text from a PDF.

    Tries ``pdfminer.six`` first (higher fidelity), then falls back to
    ``pypdf``.  Returns ``None`` if neither library is available.
    """
    # --- pdfminer.six ---
    try:
        from pdfminer.high_level import extract_text as _pdfminer_extract
        text = _pdfminer_extract(path, maxpages=50)
        if text:
            return _truncate(normalize_whitespace(text))
    except ImportError:
        pass
    except Exception as exc:
        logger.debug("search_index: pdfminer failed for %s: %s", path, exc)

    # --- pypdf fallback ---
    try:
        import pypdf  # type: ignore[import]
        pages: List[str] = []
        with open(path, "rb") as fh:
            reader = pypdf.PdfReader(fh)
            for page in reader.pages[:50]:
                page_text = page.extract_text() or ""
                pages.append(page_text)
        text = " ".join(pages)
        return _truncate(normalize_whitespace(text)) or None
    except ImportError:
        pass
    except Exception as exc:
        logger.debug("search_index: pypdf failed for %s: %s", path, exc)

    return None


def _extract_docx(path: str) -> Optional[str]:
    """Extract text from a Word document (.docx)."""
    try:
        import docx  # type: ignore[import]
        doc = docx.Document(path)
        parts: List[str] = []
        for para in doc.paragraphs:
            parts.append(para.text)
        # Also grab table cell text
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
        text = " ".join(parts)
        return _truncate(normalize_whitespace(text)) or None
    except ImportError:
        pass
    except Exception as exc:
        logger.debug("search_index: python-docx failed for %s: %s", path, exc)
    return None


def _extract_xlsx(path: str) -> Optional[str]:
    """Extract text from an Excel workbook (.xlsx/.xlsm)."""
    try:
        import openpyxl  # type: ignore[import]
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts: List[str] = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        parts.append(str(cell))
        wb.close()
        text = " ".join(parts)
        return _truncate(normalize_whitespace(text)) or None
    except ImportError:
        pass
    except Exception as exc:
        logger.debug("search_index: openpyxl failed for %s: %s", path, exc)
    return None


def _extract_xls(path: str) -> Optional[str]:
    """Extract text from a legacy Excel workbook (.xls) via xlrd."""
    try:
        import xlrd  # type: ignore[import]
        wb = xlrd.open_workbook(path)
        parts: List[str] = []
        for sheet in wb.sheets():
            for rowx in range(sheet.nrows):
                for cell in sheet.row(rowx):
                    if cell.value:
                        parts.append(str(cell.value))
        text = " ".join(parts)
        return _truncate(normalize_whitespace(text)) or None
    except ImportError:
        pass
    except Exception as exc:
        logger.debug("search_index: xlrd failed for %s: %s", path, exc)
    return None


def _extract_pptx(path: str) -> Optional[str]:
    """Extract text from a PowerPoint presentation (.pptx)."""
    try:
        from pptx import Presentation  # type: ignore[import]
        prs = Presentation(path)
        parts: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    parts.append(shape.text)
        text = " ".join(parts)
        return _truncate(normalize_whitespace(text)) or None
    except ImportError:
        pass
    except Exception as exc:
        logger.debug("search_index: python-pptx failed for %s: %s", path, exc)
    return None


def _extract_image_ocr(path: str) -> Optional[str]:
    """
    Run OCR on an image using pytesseract.

    Silently returns ``None`` if pytesseract or the Tesseract binary is
    not available.
    """
    try:
        import pytesseract  # type: ignore[import]
        from PIL import Image

        with Image.open(path) as img:
            # Convert to greyscale for better OCR accuracy
            grey = img.convert("L")
            text = pytesseract.image_to_string(grey)
        return _truncate(normalize_whitespace(text)) or None
    except ImportError:
        pass
    except Exception as exc:
        logger.debug("search_index: OCR failed for %s: %s", path, exc)
    return None


# ---------------------------------------------------------------------------
# Routing logic
# ---------------------------------------------------------------------------


def should_extract_text(filename: str, mime_type: Optional[str], file_type: str) -> bool:
    """Return True if any extractor might produce text for this file."""
    if file_type in TEXT_FILE_TYPES:
        return True
    if mime_type and mime_type.startswith("text/"):
        return True
    ext = _extension(filename)
    return ext in (
        TEXT_EXTENSIONS
        | PDF_EXTENSIONS
        | WORD_EXTENSIONS
        | EXCEL_EXTENSIONS
        | PPTX_EXTENSIONS
        | IMAGE_EXTENSIONS
    )


def extract_text_content(
    storage_path: Optional[str],
    filename: str,
    mime_type: Optional[str],
    file_type: str,
) -> Optional[str]:
    """
    Dispatch to the appropriate extractor and return normalised text, or
    ``None`` if nothing could be extracted.
    """
    if not storage_path or not os.path.exists(storage_path):
        return None

    ext = _extension(filename)

    # --- Plain text / source code ---
    if file_type in TEXT_FILE_TYPES or ext in TEXT_EXTENSIONS or (
        mime_type and mime_type.startswith("text/")
    ):
        return _extract_plain_text(storage_path)

    # --- PDF ---
    if ext in PDF_EXTENSIONS or (mime_type and "pdf" in mime_type):
        return _extract_pdf(storage_path)

    # --- Word ---
    if ext in WORD_EXTENSIONS:
        return _extract_docx(storage_path)

    # --- Excel ---
    if ext in EXCEL_EXTENSIONS:
        if ext == "xls":
            return _extract_xls(storage_path)
        return _extract_xlsx(storage_path)

    # --- PowerPoint ---
    if ext in PPTX_EXTENSIONS:
        return _extract_pptx(storage_path)

    # --- Image OCR (last resort, expensive) ---
    if ext in IMAGE_EXTENSIONS or file_type == "image":
        return _extract_image_ocr(storage_path)

    return None


# ---------------------------------------------------------------------------
# Public API (unchanged signatures for backwards-compatibility)
# ---------------------------------------------------------------------------


def build_search_document(
    storage_path: Optional[str],
    filename: str,
    mime_type: Optional[str],
    file_type: str,
) -> Optional[str]:
    """
    Top-level entry point.  Returns the extracted text or ``None``.

    Called from:
    - The upload endpoints (synchronously on the hot path for small files)
    - The background ``search_index`` job handler
    - The startup backfill task
    """
    return extract_text_content(storage_path, filename, mime_type, file_type)


def build_match_context(
    file: FileModel,
    query: str,
    path_segments: Optional[Iterable[str]] = None,
) -> Optional[str]:
    """
    Return a short text snippet around the first occurrence of *query*
    in the file's indexed content, suitable for display in search results.
    """
    normalized_query = normalize_whitespace(query).lower()
    if not normalized_query:
        return None

    haystacks = [
        ("name", file.name or ""),
        ("path", " / ".join(path_segments or [])),
        ("type", file.type or ""),
        ("mime", file.mime_type or ""),
        ("content", file.content_index or ""),
    ]

    for label, text in haystacks:
        normalized_text = normalize_whitespace(text)
        if not normalized_text:
            continue

        match_index = normalized_text.lower().find(normalized_query)
        if match_index == -1:
            continue

        if label == "content":
            start = max(0, match_index - 60)
            end = min(len(normalized_text), match_index + len(normalized_query) + 100)
            snippet = normalized_text[start:end]
            if start > 0:
                snippet = f"…{snippet}"
            if end < len(normalized_text):
                snippet = f"{snippet}…"
            return snippet[:MAX_MATCH_CONTEXT_CHARS]

        return normalized_text[:MAX_MATCH_CONTEXT_CHARS]

    return None
